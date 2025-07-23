import os
import sys
import csv
import docx
import json
import base64
import whisper
import tempfile
import pandas as pd
from PIL import Image
from moviepy import *
from cnocr import CnOcr
from PyPDF2 import PdfReader
from pptx import Presentation
from Bio.PDB import PDBParser
from bs4 import BeautifulSoup
from munch import DefaultMunch  # nested dict to object
import xml.etree.ElementTree as ET
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from llm.registry import get_llm


class SemanticRead:
    MAX_CONTENT_LENGTH = 10000  # TODO: avoid hard-coding

    def __init__(self):
        # These models are independent of the agent framework and should be deemed as a third-party tool
        # therefore, we do not enable configuration of them and assume that they are just off-the-shelf product
        # This is why we have these hard-coded
        vllm_config = {
            "type": "openai",
            "api_key_type": "ark",
            "model_name": "ep-20241217114719-kcgs9",  # doubao vision pro 32k
            "base_url": "https://ark.cn-beijing.volces.com/api/v3"
        }
        vllm_config = DefaultMunch.fromDict(vllm_config)
        self.vllm = get_llm(vllm_config)

        llm_config = {
            "type": "openai",
            "api_key_type": "ark",
            "model_name": "ep-20250212105505-5zlbx",  # doubao 1.5 pro
            "base_url": "https://ark.cn-beijing.volces.com/api/v3"
        }
        llm_config = DefaultMunch.fromDict(llm_config)
        self.llm = get_llm(llm_config)

    def extract_path_list(self, overall_task, task):
        system_prompt = '''
Basic on the context and the current specific task, please figure out the absolute file path of the file in question.
Please strictly format your answer as: FILE PATH: absolute_file_path
'''
        user_query = f'''
Overall task: {overall_task}
Current subtask: {task}
'''

        retry_time = 0
        retry_limit = 3
        file_path_list = []
        while retry_time <= retry_limit:
            retry_time += 1
            response = self.llm.get_response(
                system_prompt=system_prompt,
                user_query=user_query
            )

            lines = response.split('\n')
            for line in lines:
                if "FILE PATH: " in line:
                    file_path_list.append(line.split("FILE PATH: ")[1].strip())

            if not file_path_list:
                print(f"Retrying for the {retry_time}th time "
                      f"as file path cannot be found in the model response")
                continue
            break

        return file_path_list

    @staticmethod
    def detect_file_type(file_path):
        ext = os.path.splitext(file_path)[1].lower()
        if ext in ['.txt', '.md', ".py"]:
            return 'text'
        elif ext in ['.jpg', '.jpeg', '.apng', '.png', '.gif', '.webp', '.bmp', '.tiff', '.tif',
                     '.ico', '.dib', '.icns', '.sgi', '.j2c', '.j2k', '.jp2', '.jpc', '.jpf', '.jpx']:
            return 'image'
        elif ext in ['.mp4', '.avi', '.mov', '.mkv']:
            return 'video'
        elif ext in ['.doc', '.docx']:
            return 'word'
        elif ext in ['.xls', '.xlsx']:
            return 'excel'
        elif ext in ['.pptx', '.ppt']:
            return 'powerpoint'
        elif ext in ['.mp3', '.wav', '.aac', '.flac']:
            return 'audio'
        elif ext in ['.pdb']:
            return 'pdb'
        elif ext in ['.pdf']:
            return 'pdf'
        elif ext in ['.jsonld']:
            return 'jsonld'
        elif ext in ['.csv']:
            return 'csv'
        elif ext in ['.json']:
            return 'json'
        elif ext in ['.xml']:
            return 'xml'
        elif ext in ['.html', '.htm']:
            return 'html'
        else:
            return 'unknown'

    def process_text_file(self, file_path):
        """Read and optionally clip a text file."""
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()

        return {"type": "text", "content": content.strip()[:self.MAX_CONTENT_LENGTH]}

    def get_image_caption(self, img_bytes, ext):
        img_base64 = base64.b64encode(img_bytes).decode('utf-8')

        if ext in ['.jpg', '.jpeg']:
            content_type = 'image/jpeg'
        elif ext in ['.apng', '.png']:
            content_type = 'image/png'
        elif ext in ['.gif']:
            content_type = 'image/gif'
        elif ext in ['.webp']:
            content_type = 'image/webp'
        elif ext in ['.bmp', '.dib']:
            content_type = 'image/bmp'
        elif ext in ['.tiff', '.tif']:
            content_type = 'image/tiff'
        elif ext in ['.ico']:
            content_type = "image/x-icon"
        elif ext in ['.icns']:
            content_type = "image/icns"
        elif ext in ['.sgi']:
            content_type = "image/sgi"
        elif ext in ['.j2c', '.j2k', '.jp2', '.jpc', '.jpf', '.jpx']:
            content_type = "image/jp2"
        else:
            raise NotImplementedError

        image_url = f"data:{content_type};base64,{img_base64}"
        caption = self.vllm.get_response(
            user_query="Please generate a caption for this image.",
            image_url=image_url
        )
        return caption

    def process_image_file(self, file_path):
        """Process an image file: perform OCR and generate a caption."""
        output = CnOcr().ocr(file_path)
        extracted_text = []
        for d in output:
            text, score, position = d["text"], d["score"], d["position"]
            extracted_text.append(text.strip())
        ocr_text = " ".join(extracted_text)

        with open(file_path, 'rb') as img_file:
            img_bytes = img_file.read()
            caption = self.get_image_caption(img_bytes, os.path.splitext(file_path)[1].lower())
        return {
            "type": "image",
            "ocr_text": ocr_text.strip()[:self.MAX_CONTENT_LENGTH],
            "caption": caption.strip()[:self.MAX_CONTENT_LENGTH]
        }

    def process_pdf_file(self, file_path):
        reader = PdfReader(file_path)
        text = ""
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
        return {"type": "pdf", "content": text.strip()[:self.MAX_CONTENT_LENGTH]}

    @staticmethod
    def process_word_file(file_path):
        """Process a Word file: extract text content."""
        doc = docx.Document(file_path)
        text = "\n".join([para.text for para in doc.paragraphs])
        return {"type": "word", "content": text}

    @staticmethod
    def process_excel_file(file_path):
        """Process an Excel file: extract sheet data as dictionaries."""
        xls = pd.ExcelFile(file_path)
        sheets_data = {}
        for sheet in xls.sheet_names:
            df = pd.read_excel(xls, sheet_name=sheet)
            sheets_data[sheet] = df.to_dict(orient='records')
        return {"type": "excel", "content": sheets_data}

    @staticmethod
    def process_powerpoint_file(file_path):
        """Process a PowerPoint file: extract text from slides."""
        prs = Presentation(file_path)
        slides_data = []
        for slide in prs.slides:
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text += shape.text + "\n"
            slides_data.append({
                "slide_id": slide.slide_id,
                "text": slide_text.strip()
            })
        return {"type": "powerpoint", "slides": slides_data}

    @staticmethod
    def transcribe_audio(file_path):
        """Transcribe an audio file using a speech-to-text model (e.g., Whisper)."""
        model = whisper.load_model("base")
        result = model.transcribe(file_path, fp16=False)
        return result.get("text", "")

    def process_audio_file(self, file_path):
        """Process an audio file: perform transcription."""
        transcript = self.transcribe_audio(file_path)
        return {"type": "audio", "transcript": transcript}

    def process_video_file(self, file_path):
        """Process a video file: extract audio for transcription and a key frame for a caption."""
        clip = VideoFileClip(file_path)
        audio_temp = tempfile.NamedTemporaryFile(suffix=".wav", delete=False)
        audio_path = audio_temp.name
        clip.audio.write_audiofile(audio_path, logger=None)
        transcript = self.transcribe_audio(audio_path)
        os.remove(audio_path)

        frame_time = clip.duration / 2
        frame = clip.get_frame(frame_time)
        temp_img_file = tempfile.NamedTemporaryFile(suffix=".jpg", delete=False)
        temp_img_path = temp_img_file.name
        temp_img_file.close()
        img = Image.fromarray(frame)
        img.save(temp_img_path, format="JPEG")

        with open(temp_img_path, "rb") as f:
            img_bytes = f.read()
        caption = self.get_image_caption(img_bytes=img_bytes, ext='.jpg')
        os.remove(temp_img_path)

        return {"type": "video", "transcript": transcript, "half_duration_frame_caption": caption}

    @staticmethod
    def process_pdb_file(file_path):
        """Process a PDB file: extract basic structure summary."""
        parser = PDBParser(QUIET=True)
        structure = parser.get_structure("structure", file_path)
        summary = {}
        # Iterate over models and chains to count residues per chain
        for model in structure:
            for chain in model:
                chain_id = chain.id
                residues = list(chain.get_residues())
                summary[chain_id] = len(residues)
        return {"type": "pdb", "summary": summary}

    @staticmethod
    def process_csv_file(file_path):
        """Process a CSV file: extract its rows as a list of dictionaries."""
        with open(file_path, 'r', encoding='utf-8') as f:
            reader = csv.DictReader(f)
            rows = list(reader)
        return {"type": "csv", "content": rows}

    @staticmethod
    def process_json_file(file_path):
        """Process a JSON file: extract and return its JSON content."""
        with open(file_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        return {"type": "json", "content": data}

    @staticmethod
    def process_jsonld_file(file_path):
        """Process a JSON-LD file: extract and return its JSON content."""
        with open(file_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        return {"type": "jsonld", "content": data}

    @staticmethod
    def process_xml_file(file_path):
        """Process an XML file: extract and return its content as a string."""
        import xml.etree.ElementTree as ET
        tree = ET.parse(file_path)
        root = tree.getroot()
        xml_str = ET.tostring(root, encoding='unicode')
        return {"type": "xml", "content": xml_str}

    @staticmethod
    def process_html_file(file_path):
        """Process an HTML file: extract and return its text content."""
        with open(file_path, 'r', encoding='utf-8') as f:
            html_content = f.read()
        soup = BeautifulSoup(html_content, 'html.parser')
        # Extract text content from HTML
        text = soup.get_text(separator="\n", strip=True)
        return {"type": "html", "content": text}

    def _process(self, file_path):
        if not os.path.exists(file_path):
            raise ValueError(f"File not found: {file_path}")

        file_type = self.detect_file_type(file_path)
        if file_type == 'text':
            result = self.process_text_file(file_path)
        elif file_type == 'image':
            result = self.process_image_file(file_path)
        elif file_type == 'pdf':
            result = self.process_pdf_file(file_path)
        elif file_type == 'word':
            result = self.process_word_file(file_path)
        elif file_type == 'excel':
            result = self.process_excel_file(file_path)
        elif file_type == 'powerpoint':
            result = self.process_powerpoint_file(file_path)
        elif file_type == 'audio':
            result = self.process_audio_file(file_path)
        elif file_type == 'video':
            result = self.process_video_file(file_path)
        elif file_type == 'pdb':
            result = self.process_pdb_file(file_path)
        elif file_type == 'csv':
            result = self.process_csv_file(file_path)
        elif file_type == 'json':
            result = self.process_json_file(file_path)
        elif file_type == 'jsonld':
            result = self.process_jsonld_file(file_path)
        elif file_type == 'xml':
            result = self.process_xml_file(file_path)
        elif file_type == 'html':
            result = self.process_html_file(file_path)
        else:
            raise ValueError("Unsupported file type.")

        return result

    def process(self, file_path_list):
        result_list = []
        for file_path in file_path_list:
            result_list.append(self._process(file_path))
        return result_list
